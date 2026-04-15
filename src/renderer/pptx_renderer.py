"""
PPTX Renderer Module
Transforms designed slide plans into a valid, visually professional .pptx file
using python-pptx. Implements all layout templates with precise positioning.

Layout Templates:
  title_hero, agenda_numbered, exec_summary, two_column, bullet_clean,
  full_text, card_grid, chart_full, stat_callout, stat_grid,
  process_flow, conclusion_dark
"""

import io
import math
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree

from src.charts.chart_generator import ChartGenerator


def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Slide dimensions: 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Typography
FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_ACCENT = "Calibri Light"

# Size presets (Pt)
SZ_HERO = 52
SZ_TITLE = 36
SZ_SUBTITLE = 20
SZ_HEADING = 24
SZ_SUBHEADING = 18
SZ_BODY = 14
SZ_CAPTION = 11
SZ_STAT = 56


class PPTXRenderer:
    """Renders designed slide list into a .pptx file."""

    def __init__(self, chart_generator: ChartGenerator = None, verbose: bool = False):
        self.chart_gen = chart_generator or ChartGenerator()
        self.verbose = verbose

    def render(self, slides: list[dict], output_path: str):
        """Main entry point: render all slides and save."""
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        # Use blank layout (index 6 is blank in default template)
        blank_layout = prs.slide_layouts[6]

        for slide_data in slides:
            slide = prs.slides.add_slide(blank_layout)
            layout = slide_data.get("layout", "bullet_clean")

            if self.verbose:
                print(f"     Rendering [{slide_data['order']:2d}] {layout}: {slide_data.get('heading', '')[:40]}")

            try:
                renderer = getattr(self, f"_render_{layout}", self._render_bullet_clean)
                renderer(slide, slide_data)
            except Exception as e:
                print(f"   ⚠️  Slide {slide_data.get('order')} ({layout}) error: {e}")
                self._render_fallback(slide, slide_data)

        prs.save(output_path)

    # ── Helper: fill slide background ────────────────────────────────────────

    def _fill_bg(self, slide, hex_color: str):
        """Fill slide background with solid color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(hex_color)

    # ── Helper: add text box ─────────────────────────────────────────────────

    def _add_text(self, slide, text: str, left, top, width, height,
                  font_size: int = SZ_BODY, bold: bool = False, italic: bool = False,
                  color: str = "FFFFFF", align=PP_ALIGN.LEFT,
                  font_name: str = FONT_BODY, wrap: bool = True,
                  word_wrap: bool = True) -> Any:
        """Add a styled text box to the slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = hex_to_rgb(color)
        return txBox

    def _add_multiline_text(self, slide, lines: list[str], left, top, width, height,
                             font_size: int = SZ_BODY, color: str = "FFFFFF",
                             bold: bool = False, bullet_char: str = None,
                             line_spacing: float = 1.2) -> Any:
        """Add a text box with multiple paragraphs."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            if bullet_char:
                p.level = 0
            run = p.add_run()
            prefix = f"{bullet_char} " if bullet_char else ""
            run.text = prefix + str(line)
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = hex_to_rgb(color)

        return txBox

    # ── Helper: add shape rectangle ──────────────────────────────────────────

    def _add_rect(self, slide, left, top, width, height, fill_hex: str,
                  alpha: float = 1.0, line_color: str = None, line_width: float = 0):
        """Add a filled rectangle."""
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        if line_color:
            shape.line.color.rgb = hex_to_rgb(line_color)
            shape.line.width = Pt(line_width)
        else:
            shape.line.fill.background()
        return shape

    # ── Helper: add image from bytes ─────────────────────────────────────────

    def _add_image_bytes(self, slide, img_bytes: bytes, left, top, width, height):
        """Add image from bytes buffer."""
        if not img_bytes:
            return None
        buf = io.BytesIO(img_bytes)
        try:
            pic = slide.shapes.add_picture(buf, left, top, width, height)
            return pic
        except Exception:
            return None

    # ── Helper: slide header bar ─────────────────────────────────────────────

    def _add_header(self, slide, heading: str, slide_data: dict,
                    top_margin: float = 0.3, include_accent: bool = True):
        """Add slide title with accent strip on left."""
        accent = slide_data.get("accent_color", "F5A623")
        text_color = slide_data.get("text_color", "FFFFFF")

        if include_accent:
            # Vertical accent bar on left
            self._add_rect(slide,
                           left=Inches(0.4), top=Inches(top_margin),
                           width=Inches(0.07), height=Inches(0.65),
                           fill_hex=accent)

        self._add_text(slide, heading,
                       left=Inches(0.6 if include_accent else 0.4),
                       top=Inches(top_margin - 0.05),
                       width=Inches(12.2),
                       height=Inches(0.75),
                       font_size=SZ_TITLE, bold=True,
                       color=text_color,
                       font_name=FONT_TITLE)

    # ── Layout: Title Hero ────────────────────────────────────────────────────

    def _render_title_hero(self, slide, data: dict):
        bg = data.get("bg_color", "0D1B2A")
        accent = data.get("accent_color", "F5A623")
        primary = data.get("primary_color", "1E2761")

        self._fill_bg(slide, bg)

        # Decorative geometric shapes
        self._add_rect(slide,
                       left=Inches(0), top=Inches(0),
                       width=Inches(4.5), height=SLIDE_H,
                       fill_hex=primary)

        self._add_rect(slide,
                       left=Inches(4.2), top=Inches(2.5),
                       width=Inches(0.08), height=Inches(2.5),
                       fill_hex=accent)

        # Title
        title = data.get("heading", "Presentation Title")
        self._add_text(slide, title,
                       left=Inches(4.5), top=Inches(1.8),
                       width=Inches(8.3), height=Inches(2.5),
                       font_size=SZ_HERO, bold=True,
                       color="FFFFFF",
                       font_name=FONT_TITLE,
                       align=PP_ALIGN.LEFT)

        # Subtitle
        subtitle = data.get("subtitle", "")
        if subtitle:
            self._add_text(slide, subtitle,
                           left=Inches(4.5), top=Inches(4.5),
                           width=Inches(8.3), height=Inches(1.2),
                           font_size=SZ_SUBTITLE, bold=False,
                           color="AABCD0",
                           font_name=FONT_ACCENT,
                           align=PP_ALIGN.LEFT)

        # Left panel decoration
        self._add_text(slide, "▶",
                       left=Inches(0.5), top=Inches(3.2),
                       width=Inches(3.5), height=Inches(1.5),
                       font_size=80, bold=False,
                       color=accent,
                       align=PP_ALIGN.CENTER)

    # ── Layout: Agenda ────────────────────────────────────────────────────────

    def _render_agenda_numbered(self, slide, data: dict):
        bg = data.get("bg_color", "1E2761")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")

        self._fill_bg(slide, bg)
        self._add_header(slide, data.get("heading", "Agenda"), data, top_margin=0.3)

        items = data.get("items", [])
        n = len(items)
        cols = 2 if n > 4 else 1
        col_w = 5.8 if cols == 2 else 10.0
        start_x = [0.5, 7.0] if cols == 2 else [1.8]
        per_col = math.ceil(n / cols)

        for i, item in enumerate(items):
            col = i // per_col
            row = i % per_col
            x = start_x[col] if col < len(start_x) else start_x[-1]
            y = 1.4 + row * 1.0

            # Number circle
            self._add_rect(slide,
                           left=Inches(x), top=Inches(y),
                           width=Inches(0.5), height=Inches(0.5),
                           fill_hex=accent)
            self._add_text(slide, str(i + 1),
                           left=Inches(x), top=Inches(y),
                           width=Inches(0.5), height=Inches(0.5),
                           font_size=14, bold=True,
                           color="000000" if accent == "F5A623" else "FFFFFF",
                           align=PP_ALIGN.CENTER)

            # Item text
            self._add_text(slide, item,
                           left=Inches(x + 0.65), top=Inches(y + 0.05),
                           width=Inches(col_w - 0.7), height=Inches(0.5),
                           font_size=16, bold=False,
                           color=text_color)

    # ── Layout: Executive Summary ─────────────────────────────────────────────

    def _render_exec_summary(self, slide, data: dict):
        bg = data.get("bg_color", "1E2761")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")
        secondary = data.get("secondary_color", "7EC8E3")

        self._fill_bg(slide, bg)

        # Dark left panel
        self._add_rect(slide,
                       left=Inches(0), top=Inches(0),
                       width=Inches(3.2), height=SLIDE_H,
                       fill_hex="00000033" if bg != "0D1B2A" else "0D1B2A")

        # Left panel label
        self._add_text(slide, "EXEC\nSUMMARY",
                       left=Inches(0.2), top=Inches(2.8),
                       width=Inches(2.8), height=Inches(2.0),
                       font_size=20, bold=True,
                       color=accent,
                       align=PP_ALIGN.CENTER)

        self._add_text(slide, data.get("heading", "Executive Summary"),
                       left=Inches(3.4), top=Inches(0.3),
                       width=Inches(9.5), height=Inches(0.8),
                       font_size=SZ_TITLE, bold=True,
                       color=text_color,
                       font_name=FONT_TITLE)

        bullets = data.get("bullets", [])
        for i, bullet in enumerate(bullets[:5]):
            y = 1.3 + i * 1.0
            # Accent dot
            self._add_rect(slide,
                           left=Inches(3.4), top=Inches(y + 0.15),
                           width=Inches(0.12), height=Inches(0.12),
                           fill_hex=accent)
            self._add_text(slide, bullet,
                           left=Inches(3.7), top=Inches(y),
                           width=Inches(9.2), height=Inches(0.9),
                           font_size=14, color=text_color)

    # ── Layout: Two Column ────────────────────────────────────────────────────

    def _render_two_column(self, slide, data: dict):
        bg = data.get("bg_color", "1E2761")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")
        card_bg = data.get("secondary_color", "243447")

        self._fill_bg(slide, bg)
        self._add_header(slide, data.get("heading", ""), data)

        bullets = data.get("bullets", [])
        half = math.ceil(len(bullets) / 2)
        left_bullets = bullets[:half]
        right_bullets = bullets[half:]

        for col_i, col_bullets in enumerate([left_bullets, right_bullets]):
            x = 0.4 + col_i * 6.5
            for j, bullet in enumerate(col_bullets[:4]):
                y = 1.5 + j * 1.1
                # Card background
                self._add_rect(slide,
                               left=Inches(x), top=Inches(y),
                               width=Inches(6.0), height=Inches(0.9),
                               fill_hex=card_bg if data.get("is_dark_bg") else "FFFFFF")
                # Accent left border
                self._add_rect(slide,
                               left=Inches(x), top=Inches(y),
                               width=Inches(0.08), height=Inches(0.9),
                               fill_hex=accent)
                self._add_text(slide, bullet,
                               left=Inches(x + 0.2), top=Inches(y + 0.1),
                               width=Inches(5.6), height=Inches(0.7),
                               font_size=13, color=text_color)

    # ── Layout: Bullet Clean ──────────────────────────────────────────────────

    def _render_bullet_clean(self, slide, data: dict):
        bg = data.get("bg_color", "1E2761")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")

        self._fill_bg(slide, bg)
        self._add_header(slide, data.get("heading", ""), data)

        bullets = data.get("bullets", [])
        paragraphs = data.get("paragraphs", [])

        y_start = 1.4
        # Leading paragraph
        if paragraphs:
            self._add_text(slide, paragraphs[0],
                           left=Inches(0.5), top=Inches(y_start),
                           width=Inches(12.3), height=Inches(0.9),
                           font_size=14, color=text_color,
                           font_name=FONT_ACCENT, italic=True)
            y_start = 2.4

        for i, bullet in enumerate(bullets[:6]):
            y = y_start + i * 0.85
            if y > 6.8:
                break
            # Bullet indicator
            self._add_rect(slide,
                           left=Inches(0.5), top=Inches(y + 0.2),
                           width=Inches(0.12), height=Inches(0.12),
                           fill_hex=accent)
            self._add_text(slide, bullet,
                           left=Inches(0.8), top=Inches(y),
                           width=Inches(12.0), height=Inches(0.75),
                           font_size=15, color=text_color)

    # ── Layout: Full Text ─────────────────────────────────────────────────────

    def _render_full_text(self, slide, data: dict):
        bg = data.get("bg_color", "1E2761")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")

        self._fill_bg(slide, bg)
        self._add_header(slide, data.get("heading", ""), data)

        paragraphs = data.get("paragraphs", [])
        bullets = data.get("bullets", [])

        y = 1.5
        for para in paragraphs[:2]:
            self._add_text(slide, para,
                           left=Inches(0.5), top=Inches(y),
                           width=Inches(12.3), height=Inches(1.5),
                           font_size=15, color=text_color)
            y += 1.7

        for bullet in bullets[:3]:
            if y > 6.5:
                break
            self._add_rect(slide,
                           left=Inches(0.5), top=Inches(y + 0.22),
                           width=Inches(0.1), height=Inches(0.1),
                           fill_hex=accent)
            self._add_text(slide, bullet,
                           left=Inches(0.75), top=Inches(y),
                           width=Inches(12.0), height=Inches(0.7),
                           font_size=14, color=text_color)
            y += 0.85

    # ── Layout: Card Grid ─────────────────────────────────────────────────────

    def _render_card_grid(self, slide, data: dict):
        bg = data.get("bg_color", "1E2761")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")
        card_bg = "243447" if data.get("is_dark_bg") else "E8F4F8"

        self._fill_bg(slide, bg)
        self._add_header(slide, data.get("heading", ""), data)

        bullets = data.get("bullets", [])[:6]
        n = len(bullets)
        cols = 3 if n >= 4 else 2
        rows = math.ceil(n / cols)
        card_w = 12.0 / cols
        card_h = 4.8 / rows

        for i, bullet in enumerate(bullets):
            col = i % cols
            row = i // cols
            x = 0.55 + col * (card_w + 0.1)
            y = 1.5 + row * (card_h + 0.15)

            # Card
            self._add_rect(slide,
                           left=Inches(x), top=Inches(y),
                           width=Inches(card_w - 0.1), height=Inches(card_h - 0.1),
                           fill_hex=card_bg,
                           line_color=accent, line_width=0.5)

            # Top accent bar
            self._add_rect(slide,
                           left=Inches(x), top=Inches(y),
                           width=Inches(card_w - 0.1), height=Inches(0.08),
                           fill_hex=accent)

            # Number
            self._add_text(slide, str(i + 1),
                           left=Inches(x + 0.1), top=Inches(y + 0.15),
                           width=Inches(0.4), height=Inches(0.4),
                           font_size=18, bold=True, color=accent)

            # Text
            self._add_text(slide, bullet,
                           left=Inches(x + 0.1), top=Inches(y + 0.5),
                           width=Inches(card_w - 0.3), height=Inches(card_h - 0.7),
                           font_size=12, color=text_color)

    # ── Layout: Chart Full ────────────────────────────────────────────────────

    def _render_chart_full(self, slide, data: dict):
        bg = data.get("bg_color", "1E2761")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")

        self._fill_bg(slide, bg)
        self._add_header(slide, data.get("heading", "Data Overview"), data)

        table = data.get("table")
        chart_type = data.get("chart_type", "bar")

        if table:
            theme_name = "dark" if data.get("is_dark_bg") else "light"
            self.chart_gen.set_theme(theme_name)
            chart_bytes = self.chart_gen.generate(table, chart_type,
                                                   width_in=10, height_in=5.2)
            if chart_bytes:
                self._add_image_bytes(slide, chart_bytes,
                                      left=Inches(1.2), top=Inches(1.3),
                                      width=Inches(10.9), height=Inches(5.7))
                return

        # Fallback: render table as text
        if table:
            headers = table.get("headers", [])
            rows = table.get("rows", [])[:8]
            self._render_table_text(slide, headers, rows, text_color, bg)

    def _render_table_text(self, slide, headers, rows, text_color, bg):
        """Simple text table fallback."""
        col_w = 12.0 / max(len(headers), 1)
        y = 1.4
        # Header row
        for j, h in enumerate(headers[:6]):
            self._add_text(slide, h,
                           left=Inches(0.5 + j * col_w), top=Inches(y),
                           width=Inches(col_w - 0.1), height=Inches(0.5),
                           font_size=13, bold=True, color=text_color)
        y += 0.6
        for row in rows:
            for j, cell in enumerate(row[:6]):
                self._add_text(slide, str(cell),
                               left=Inches(0.5 + j * col_w), top=Inches(y),
                               width=Inches(col_w - 0.1), height=Inches(0.45),
                               font_size=12, color=text_color)
            y += 0.5
            if y > 7.0:
                break

    # ── Layout: Stat Callout ──────────────────────────────────────────────────

    def _render_stat_callout(self, slide, data: dict):
        bg = data.get("bg_color", "0D1B2A")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")

        self._fill_bg(slide, bg)
        self._add_header(slide, data.get("heading", "Key Metrics"), data)

        stats = data.get("stats", [])[:3]
        n = len(stats)
        w = 12.0 / max(n, 1)

        for i, stat in enumerate(stats):
            x = 0.65 + i * w
            y = 2.0
            val = stat.get("value", "—")
            ctx = stat.get("context", "")[:60]

            # Value card
            self._add_rect(slide,
                           left=Inches(x + 0.1), top=Inches(y),
                           width=Inches(w - 0.3), height=Inches(3.5),
                           fill_hex=data.get("primary_color", "1E2761"))

            # Large stat value
            self._add_text(slide, val,
                           left=Inches(x + 0.1), top=Inches(y + 0.4),
                           width=Inches(w - 0.3), height=Inches(1.8),
                           font_size=SZ_STAT, bold=True,
                           color=accent, align=PP_ALIGN.CENTER)

            # Context label
            self._add_text(slide, ctx,
                           left=Inches(x + 0.15), top=Inches(y + 2.2),
                           width=Inches(w - 0.4), height=Inches(1.1),
                           font_size=12, color=text_color,
                           align=PP_ALIGN.CENTER)

    # ── Layout: Stat Grid ─────────────────────────────────────────────────────

    def _render_stat_grid(self, slide, data: dict):
        """Grid of stat cards for 4-6 stats."""
        bg = data.get("bg_color", "0D1B2A")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")
        primary = data.get("primary_color", "1E2761")

        self._fill_bg(slide, bg)
        self._add_header(slide, data.get("heading", "Key Metrics"), data)

        stats = data.get("stats", [])[:6]
        n = len(stats)
        cols = 3
        rows = math.ceil(n / cols)
        card_w = 3.9
        card_h = 2.3

        for i, stat in enumerate(stats):
            col = i % cols
            row = i // cols
            x = 0.4 + col * (card_w + 0.2)
            y = 1.5 + row * (card_h + 0.2)
            val = stat.get("value", "—")
            ctx = stat.get("context", "")[:50]

            self._add_rect(slide,
                           left=Inches(x), top=Inches(y),
                           width=Inches(card_w), height=Inches(card_h),
                           fill_hex=primary)

            self._add_rect(slide,
                           left=Inches(x), top=Inches(y),
                           width=Inches(card_w), height=Inches(0.07),
                           fill_hex=accent)

            self._add_text(slide, val,
                           left=Inches(x), top=Inches(y + 0.2),
                           width=Inches(card_w), height=Inches(1.2),
                           font_size=40, bold=True,
                           color=accent, align=PP_ALIGN.CENTER)

            self._add_text(slide, ctx,
                           left=Inches(x + 0.1), top=Inches(y + 1.4),
                           width=Inches(card_w - 0.2), height=Inches(0.8),
                           font_size=11, color=text_color,
                           align=PP_ALIGN.CENTER)

    # ── Layout: Process Flow ──────────────────────────────────────────────────

    def _render_process_flow(self, slide, data: dict):
        bg = data.get("bg_color", "1E2761")
        accent = data.get("accent_color", "F5A623")
        text_color = data.get("text_color", "FFFFFF")
        secondary = data.get("secondary_color", "7EC8E3")

        self._fill_bg(slide, bg)
        self._add_header(slide, data.get("heading", ""), data)

        steps = (data.get("bullets") or data.get("paragraphs", []))[:5]
        n = len(steps)
        step_w = 11.5 / max(n, 1)
        y = 2.5
        arrow_color = secondary

        for i, step in enumerate(steps):
            x = 0.8 + i * step_w

            # Step circle
            circle = slide.shapes.add_shape(9, Inches(x), Inches(y),
                                             Inches(0.9), Inches(0.9))
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex_to_rgb(accent)
            circle.line.fill.background()

            # Step number
            self._add_text(slide, str(i + 1),
                           left=Inches(x), top=Inches(y),
                           width=Inches(0.9), height=Inches(0.9),
                           font_size=18, bold=True,
                           color="000000" if accent in ("F5A623", "FFD700") else "FFFFFF",
                           align=PP_ALIGN.CENTER)

            # Arrow (except last)
            if i < n - 1:
                arr_x = x + 0.95
                self._add_text(slide, "→",
                               left=Inches(arr_x), top=Inches(y + 0.1),
                               width=Inches(step_w - 0.95), height=Inches(0.7),
                               font_size=20, color=arrow_color,
                               align=PP_ALIGN.CENTER)

            # Step label
            self._add_text(slide, step,
                           left=Inches(x - 0.2), top=Inches(y + 1.0),
                           width=Inches(step_w + 0.1), height=Inches(2.5),
                           font_size=12, color=text_color,
                           align=PP_ALIGN.CENTER)

    # ── Layout: Conclusion ────────────────────────────────────────────────────

    def _render_conclusion_dark(self, slide, data: dict):
        bg = data.get("bg_color", "0D1B2A")
        accent = data.get("accent_color", "F5A623")
        primary = data.get("primary_color", "1E2761")
        text_color = data.get("text_color", "FFFFFF")

        self._fill_bg(slide, bg)

        # Right decorative panel
        self._add_rect(slide,
                       left=Inches(9.3), top=Inches(0),
                       width=Inches(4.03), height=SLIDE_H,
                       fill_hex=primary)

        # Checkmark / star on right panel
        self._add_text(slide, "★",
                       left=Inches(9.3), top=Inches(2.5),
                       width=Inches(4.03), height=Inches(2.0),
                       font_size=80, color=accent,
                       align=PP_ALIGN.CENTER)

        # Title
        self._add_text(slide, data.get("heading", "Key Takeaways"),
                       left=Inches(0.5), top=Inches(0.4),
                       width=Inches(8.5), height=Inches(0.85),
                       font_size=SZ_TITLE, bold=True, color=accent,
                       font_name=FONT_TITLE)

        bullets = data.get("bullets", [])
        for i, bullet in enumerate(bullets[:5]):
            y = 1.5 + i * 0.98
            # Check icon
            self._add_text(slide, "✓",
                           left=Inches(0.5), top=Inches(y),
                           width=Inches(0.4), height=Inches(0.7),
                           font_size=16, bold=True, color=accent)
            self._add_text(slide, bullet,
                           left=Inches(1.0), top=Inches(y),
                           width=Inches(8.0), height=Inches(0.85),
                           font_size=15, color=text_color)

    # ── Fallback ──────────────────────────────────────────────────────────────

    def _render_fallback(self, slide, data: dict):
        bg = data.get("bg_color", "1E2761")
        text_color = data.get("text_color", "FFFFFF")
        self._fill_bg(slide, bg)
        self._add_text(slide, data.get("heading", "Slide"),
                       left=Inches(0.5), top=Inches(0.4),
                       width=Inches(12.3), height=Inches(0.8),
                       font_size=SZ_TITLE, bold=True, color=text_color)

        content = " ".join(data.get("bullets", []) or data.get("paragraphs", []))
        if content:
            self._add_text(slide, content,
                           left=Inches(0.5), top=Inches(1.5),
                           width=Inches(12.3), height=Inches(5.5),
                           font_size=SZ_BODY, color=text_color)
